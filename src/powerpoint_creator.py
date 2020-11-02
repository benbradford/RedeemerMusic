from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json

class PowerpointCreator:
    def __init__(self, songs, out_file):
        self._out_file = out_file
        self._songs = songs

    def create(self):
        presentation = self._create_empty_presentation()

        for song in self._songs:
            self._add_slides_for_song(presentation, song)
            self._add_blank_page(presentation)

        outfile = "bin/{}".format(self._out_file)
        presentation.save(outfile)
        return outfile

    def _create_empty_presentation(self):
        empty = Presentation()
        empty.slide_width = Inches(14)
        blank_slide_layout = empty.slide_layouts[6]
        return empty

    def _add_slides_for_song(self, presentation, song):
        paginated_lyrics = self._paginate_lyrics(song.get_ppt_file())
        num_slides_with_lyrics = len(paginated_lyrics)

        for i in range(num_slides_with_lyrics):
            self._assemble_page(paginated_lyrics[i], presentation, i == num_slides_with_lyrics-1)

    def _assemble_page(self, lyrics_on_page, presentation, is_last_slide):
        num_lines = lyrics_on_page.count('\n') + 1
        slide = self._create_slide(presentation)
        paragraph = self._create_paragraph(slide, num_lines)
        self._add_lyrics_to_paragraph(paragraph, lyrics_on_page)
        if is_last_slide:
            self._add_footer(slide, 'CCLI Licence No. 640402')

    def _add_blank_page(self, presentation):
        self._create_slide(presentation)

    def _paginate_lyrics(self, song_file):
        # '\n' denotes a line seperator
        paginated_lyrics = []
        file = open(song_file, "r")

        next_line = file.readline()
        while next_line != '':
            self._append_lyrics_block(next_line, file, paginated_lyrics)
            next_line = self._skip_blank_lines(file)

        return paginated_lyrics

    def _append_lyrics_block(self, accumulated, file, paginated_lyrics):
        next = file.readline()
        while next != '\n' and next != '' and len(next) > 2:
            accumulated = accumulated + next
            next = file.readline()
        paginated_lyrics.append(accumulated.replace('\r', ''))

    def _skip_blank_lines(self, file):
        accumulated = file.readline()
        while accumulated == '\n' or accumulated == '\r' or accumulated == ' ':
            accumulated = file.readline()
        return accumulated

    def _create_slide(self, presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        return slide

    def _create_paragraph(self, slide, num_lines):
        width = height = Inches(1)
        left = Inches(6.5)
        top = Inches(self._get_top_margin_in_inches(num_lines))
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        return tf.add_paragraph()

    def _add_lyrics_to_paragraph(self, p, lyrics):
        p.text = lyrics
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def _get_top_margin_in_inches(self, num_lines):
        lines_to_inches = {
             1: 2.75,
             2: 2.5,
             3: 2.25,
             4: 2,
             5: 1.75,
             6: 1.5,
             7: 1.25,
             8: 0.7,
             9: 0.3
        }
        try:
            return lines_to_inches[num_lines]
        except:
            raise Exception("Exceeded max allowed lines on a page {}/{}".format(num_lines, 9))

    def _add_footer(self, slide, footer):
        width = height = Inches(1)
        left = Inches(13.0)
        top = Inches(6.9)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.RIGHT
