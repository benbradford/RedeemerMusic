from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json


presentation = Presentation()

def _create_empty_presentation():
    empty = Presentation()
    empty.slide_width = Inches(14)
    blank_slide_layout = empty.slide_layouts[6]
    return empty

def _create_slide(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    return slide

def _get_top_margin_in_inches(num_lines):
    lines_to_inches = {
         1: 2.75,
         2: 2.5,
         3: 2.25,
         4: 2,
         5: 1.75,
         6: 1.5,
         7: 1.25,
         8: 0.7,
         9: 0.3,
         10: 0
    }
    try:
        return lines_to_inches[num_lines]
    except:
        raise Exception("Exceeded max allowed lines on a page {}/{}".format(num_lines, 10))

def _create_paragraph(slide, num_lines):
    width = height = Inches(1)
    left = Inches(6.5)
    top = Inches(_get_top_margin_in_inches(num_lines))
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    return tf.add_paragraph()

def _add_lyrics_to_paragraph(p, lyrics):
    text = ""
    for line in lyrics:
        text += line + '\n'
    p.text = text
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def _add_footer(slide, footer):
    width = height = Inches(1)
    left = Inches(13.0)
    top = Inches(6.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = footer
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.RIGHT

def _assemble_page(lyrics_on_page, presentation, is_last_slide):
    num_lines = len(lyrics_on_page) + 1
    slide = _create_slide(presentation)
    paragraph = _create_paragraph(slide, num_lines)
    _add_lyrics_to_paragraph(paragraph, lyrics_on_page)
    if is_last_slide:
        _add_footer(slide, 'CCLI Licence No. 640402')

def _add_slides_for_song(presentation, song):
    num_slides_with_lyrics = len(song)
    for i in range(num_slides_with_lyrics):
        _assemble_page(song[i], presentation, i == num_slides_with_lyrics-1)

def _create_slide(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    return slide

def _add_blank_page(presentation):
    _create_slide(presentation)

def create_pp(songs, filename):
    presentation = _create_empty_presentation()
    for song in songs:
        _add_slides_for_song(presentation, song['slides'])
        _add_blank_page(presentation)
    presentation.save(filename)
    return filename
