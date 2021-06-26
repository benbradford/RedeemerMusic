from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json

class SlidesHelper:
    def __init__(self, data_retriever):
        self._data_retriever = data_retriever

    def create_powerpoint(self, service, out_file):
        slides = self._get_slides_from_service(service)
        presentation = self._create_empty_presentation()

        for slide in slides:
            pages = self._break_down_slides_to_pages(slide)
            self._add_slides_for_song(presentation, pages)
            self._add_blank_page(presentation)

        presentation.save(out_file)
        return out_file

    def _break_down_slides_to_pages(self, lyrics):
        split = lyrics.split('\n')
        pages = []
        curr = ""
        for s in split:
            if len(s) > 1:
                curr += s + '\n'
            else:
                if len(curr) > 1:
                    pages.append(curr)
                    print curr
                    curr = ""

        if len(curr) > 1:
            print curr
            pages.append(curr + '\n')
        return pages

    def _get_slides_from_service(self, service):
        slides = []
        for index in [1,2,3,4,5, 6]:
            song_key = "song" + str(index)
            if song_key in service and service[song_key] != "":
                s = self._data_retriever.get_slide(service[song_key])
                slides.append(s)

        return slides

    def _create_empty_presentation(self):
        empty = Presentation()
        empty.slide_width = Inches(14)
        blank_slide_layout = empty.slide_layouts[6]
        return empty

    def _add_slides_for_song(self, presentation, slide):
        num_slides_with_lyrics = len(slide)

        for i in range(num_slides_with_lyrics):
            self._assemble_page(slide[i], presentation, i == num_slides_with_lyrics-1)

    def _assemble_page(self, lyrics_on_page, presentation, is_last_slide):
        num_lines = lyrics_on_page.count('\n') + 1
        slide = self._create_slide(presentation)
        paragraph = self._create_paragraph(slide, num_lines)
        self._add_lyrics_to_paragraph(paragraph, lyrics_on_page)
        if is_last_slide:
            self._add_footer(slide, 'CCLI Licence No. 640402')

    def _add_blank_page(self, presentation):
        self._create_slide(presentation)

    def _create_slide(self, presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        return slide

    def _create_paragraph(self, slide, num_lines):
        width = height = Inches(1)
        left = Inches(6.5)
        top = Inches(self._get_top_margin_in_inches(num_lines))
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        return tf.add_paragraph()

    def _add_lyrics_to_paragraph(self, p, lyrics):
        p.text = lyrics.replace('\r', '')
        p.font.size = Pt(54)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def _get_top_margin_in_inches(self, num_lines):
        lines_to_inches = {
             1: 2.75,
             2: 2.4,
             3: 2.1,
             4: 1.7,
             5: 1.4,
             6: 1.0,
             7: 0.4,
             8: 0,
             9: 0
        }
        try:
            return lines_to_inches[num_lines]
        except:
            raise Exception("Exceeded max allowed lines on a page {}/{}".format(num_lines, 9))

    def _add_footer(self, slide, footer):
        width = height = Inches(1)
        left = Inches(13.0)
        top = Inches(6.9)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.RIGHT
